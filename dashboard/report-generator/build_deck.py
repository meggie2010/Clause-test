"""
PowerPoint Deck Builder — Fortnightly Client Spending Report

Assembles screenshots from Looker Studio into a professional
PowerPoint presentation with branded slides.

Usage:
    python build_deck.py
    python build_deck.py --date-range "Jan 1 - Jan 15, 2026"
    python build_deck.py --client "Sarah M."
"""

import argparse
import json
from datetime import datetime, timedelta
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def load_config():
    config_path = Path(__file__).parent / "config.json"
    with open(config_path) as f:
        return json.load(f)


def hex_to_rgb(hex_color):
    """Convert '#2b6777' to RGBColor."""
    h = hex_color.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_background(slide, color):
    """Set the slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_fill(slide, left, top, width, height, color):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 font_color=None, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    """Add a text box with styling."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if font_color:
        p.font.color.rgb = font_color
    return txbox


def get_default_date_range():
    """Generate a fortnightly date range ending today."""
    end = datetime.now()
    start = end - timedelta(days=14)
    return f"{start.strftime('%b %d')} - {end.strftime('%b %d, %Y')}"


def build_title_slide(prs, config, date_range, client_name):
    """Slide 1: Title slide with branding."""
    colors = config["brand_colors"]
    primary = hex_to_rgb(colors["primary"])
    accent = hex_to_rgb(colors["accent"])
    dark = hex_to_rgb(colors["dark"])
    white = RGBColor(0xFF, 0xFF, 0xFF)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_background(slide, primary)

    # Accent bar at top
    add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), accent)

    # Business name
    add_text_box(
        slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.8),
        config["business_name"],
        font_size=20, font_color=accent, bold=False
    )

    # Main title
    add_text_box(
        slide, Inches(0.8), Inches(2.2), Inches(11), Inches(1.2),
        "Financial Review",
        font_size=44, font_color=white, bold=True
    )

    # Date range
    add_text_box(
        slide, Inches(0.8), Inches(3.5), Inches(11), Inches(0.6),
        date_range,
        font_size=24, font_color=RGBColor(0xC8, 0xD8, 0xE4)
    )

    # Client name
    if client_name:
        add_text_box(
            slide, Inches(0.8), Inches(4.3), Inches(11), Inches(0.5),
            f"Prepared for {client_name}",
            font_size=18, font_color=RGBColor(0xC8, 0xD8, 0xE4)
        )

    # Consultant info at bottom
    add_text_box(
        slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.4),
        f"Prepared by {config['consultant_name']}",
        font_size=14, font_color=RGBColor(0x8A, 0xA8, 0xB4)
    )

    # Bottom accent bar
    add_shape_fill(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), accent)


def build_section_slide(prs, config, title, subtitle=""):
    """Section divider slide."""
    colors = config["brand_colors"]
    primary = hex_to_rgb(colors["primary"])
    accent = hex_to_rgb(colors["accent"])
    white = RGBColor(0xFF, 0xFF, 0xFF)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, primary)
    add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), accent)

    add_text_box(
        slide, Inches(0.8), Inches(2.8), Inches(11), Inches(1.0),
        title,
        font_size=36, font_color=white, bold=True
    )

    if subtitle:
        add_text_box(
            slide, Inches(0.8), Inches(3.9), Inches(9), Inches(0.6),
            subtitle,
            font_size=18, font_color=RGBColor(0xC8, 0xD8, 0xE4)
        )

    add_shape_fill(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), accent)


def build_screenshot_slide(prs, config, title, subtitle, screenshot_path):
    """Slide with a Looker Studio screenshot."""
    colors = config["brand_colors"]
    primary = hex_to_rgb(colors["primary"])
    accent = hex_to_rgb(colors["accent"])
    dark = hex_to_rgb(colors["dark"])
    text_color = hex_to_rgb(colors["text"])
    light = hex_to_rgb(colors["light"])
    white = RGBColor(0xFF, 0xFF, 0xFF)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, light)

    # Top bar
    add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), primary)
    add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), accent)

    # Title on the bar
    add_text_box(
        slide, Inches(0.6), Inches(0.2), Inches(11), Inches(0.7),
        title,
        font_size=24, font_color=white, bold=True
    )

    # Subtitle
    if subtitle:
        add_text_box(
            slide, Inches(0.6), Inches(1.15), Inches(11), Inches(0.4),
            subtitle,
            font_size=14, font_color=RGBColor(0x64, 0x74, 0x8B)
        )
        img_top = Inches(1.6)
    else:
        img_top = Inches(1.2)

    # Add screenshot image
    if screenshot_path and Path(screenshot_path).exists():
        img_left = Inches(0.4)
        img_width = Inches(12.5)
        img_height = Inches(5.6) if subtitle else Inches(6.0)
        slide.shapes.add_picture(
            str(screenshot_path), img_left, img_top, img_width, img_height
        )
    else:
        # Placeholder if screenshot is missing
        placeholder = add_shape_fill(
            slide, Inches(0.6), img_top, Inches(12.1), Inches(5.5),
            RGBColor(0xE2, 0xE8, 0xF0)
        )
        add_text_box(
            slide, Inches(3), Inches(3.5), Inches(7), Inches(1),
            f"Screenshot not found: {screenshot_path or 'N/A'}",
            font_size=16, font_color=RGBColor(0x94, 0xA3, 0xB8),
            alignment=PP_ALIGN.CENTER
        )

    # Bottom accent line
    add_shape_fill(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), accent)


def build_insights_slide(prs, config):
    """Slide for key insights and next steps (placeholder for AI later)."""
    colors = config["brand_colors"]
    primary = hex_to_rgb(colors["primary"])
    accent = hex_to_rgb(colors["accent"])
    text_color = hex_to_rgb(colors["text"])
    light = hex_to_rgb(colors["light"])
    white = RGBColor(0xFF, 0xFF, 0xFF)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, light)

    # Top bar
    add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), primary)
    add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), accent)

    add_text_box(
        slide, Inches(0.6), Inches(0.2), Inches(11), Inches(0.7),
        "Key Insights & Next Steps",
        font_size=24, font_color=white, bold=True
    )

    # Insights section
    add_text_box(
        slide, Inches(0.6), Inches(1.3), Inches(5.5), Inches(0.5),
        "Insights",
        font_size=20, font_color=primary, bold=True
    )

    insights_placeholder = [
        "Your top spending category this period was [Category] at $[Amount]",
        "You are [X%] through your monthly budget with [Y%] of the month remaining",
        "Spending in [Category] increased [X%] compared to last period",
    ]

    for i, insight in enumerate(insights_placeholder):
        y_pos = Inches(1.9 + (i * 0.55))
        add_text_box(
            slide, Inches(0.9), y_pos, Inches(5.5), Inches(0.5),
            f"   {insight}",
            font_size=14, font_color=text_color
        )

    # Next steps section
    add_text_box(
        slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.5),
        "Recommended Actions",
        font_size=20, font_color=primary, bold=True
    )

    actions_placeholder = [
        "Review [Category] spending for savings opportunities",
        "Consider adjusting budget for [Category] based on trends",
        "Schedule check-in to discuss [specific topic]",
    ]

    for i, action in enumerate(actions_placeholder):
        y_pos = Inches(1.9 + (i * 0.55))
        add_text_box(
            slide, Inches(7.3), y_pos, Inches(5.5), Inches(0.5),
            f"   {action}",
            font_size=14, font_color=text_color
        )

    # Bottom accent line
    add_shape_fill(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), accent)


def build_closing_slide(prs, config):
    """Final slide with contact info."""
    colors = config["brand_colors"]
    primary = hex_to_rgb(colors["primary"])
    accent = hex_to_rgb(colors["accent"])
    white = RGBColor(0xFF, 0xFF, 0xFF)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, primary)
    add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), accent)

    add_text_box(
        slide, Inches(0.8), Inches(2.2), Inches(11), Inches(0.8),
        "Thank You",
        font_size=40, font_color=white, bold=True,
        alignment=PP_ALIGN.CENTER
    )

    add_text_box(
        slide, Inches(0.8), Inches(3.3), Inches(11), Inches(0.6),
        f"Questions? Reach out anytime.",
        font_size=20, font_color=RGBColor(0xC8, 0xD8, 0xE4),
        alignment=PP_ALIGN.CENTER
    )

    add_text_box(
        slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.5),
        config["consultant_name"],
        font_size=22, font_color=accent, bold=True,
        alignment=PP_ALIGN.CENTER
    )

    add_text_box(
        slide, Inches(0.8), Inches(5.1), Inches(11), Inches(0.4),
        config["business_name"],
        font_size=16, font_color=RGBColor(0xC8, 0xD8, 0xE4),
        alignment=PP_ALIGN.CENTER
    )

    add_shape_fill(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), accent)


def build_report(config, date_range=None, client_name=None):
    """Build the full PowerPoint report."""
    if not date_range:
        date_range = get_default_date_range()

    if not client_name:
        client_name = config["client_name"]

    screenshots_dir = Path(__file__).parent / "screenshots"
    output_dir = Path(__file__).parent / config["output_dir"]
    output_dir.mkdir(exist_ok=True)

    # Widescreen 16:9 presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    print("Building title slide...")
    build_title_slide(prs, config, date_range, client_name)

    # Slides 2-5: Dashboard screenshots
    for report_page in config["report_pages"]:
        name = report_page["name"]
        label = report_page["label"]
        description = report_page.get("description", "")
        screenshot_path = screenshots_dir / f"{name}.png"

        print(f"Building slide: {label}...")
        build_screenshot_slide(prs, config, label, description, screenshot_path)

    # Slide 6: Insights
    print("Building insights slide...")
    build_insights_slide(prs, config)

    # Slide 7: Closing
    print("Building closing slide...")
    build_closing_slide(prs, config)

    # Save
    safe_client = client_name.replace(" ", "_").replace(".", "")
    date_stamp = datetime.now().strftime("%Y-%m-%d")
    filename = f"Financial_Review_{safe_client}_{date_stamp}.pptx"
    output_path = output_dir / filename

    prs.save(str(output_path))
    print(f"\nReport saved: {output_path}")
    return output_path


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Build a fortnightly financial review PowerPoint deck"
    )
    parser.add_argument(
        "--date-range",
        help='Date range for the report (e.g., "Jan 1 - Jan 15, 2026")'
    )
    parser.add_argument(
        "--client",
        help="Client name (overrides config.json)"
    )
    args = parser.parse_args()

    config = load_config()
    build_report(config, date_range=args.date_range, client_name=args.client)
